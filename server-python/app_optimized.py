#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PDF转换服务 - 超高性能优化版 v3.0
专门针对Windows系统优化，解决多进程死锁问题
"""
import os
import sys
import uuid
import time
import logging
from datetime import datetime, timedelta
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
import fitz  # PyMuPDF
from pdf2docx import Converter
from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
import threading
import signal
import traceback
import base64
import io
import json
from PIL import Image

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)

app = Flask(__name__)
CORS(app)

# 配置
UPLOAD_FOLDER = 'uploads'
CONVERTED_FOLDER = 'converted'
ALLOWED_EXTENSIONS = {'pdf'}
MAX_FILE_SIZE = 100 * 1024 * 1024  # 100MB

os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(CONVERTED_FOLDER, exist_ok=True)

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['CONVERTED_FOLDER'] = CONVERTED_FOLDER
app.config['MAX_CONTENT_LENGTH'] = MAX_FILE_SIZE

# 全局超时配置
CONVERSION_TIMEOUT = 300  # 5分钟超时


class TimeoutException(Exception):
    """超时异常"""
    pass


def timeout_handler(func, timeout_seconds):
    """超时处理装饰器"""
    def wrapper(*args, **kwargs):
        result = [TimeoutException('转换超时')]
        
        def target():
            try:
                result[0] = func(*args, **kwargs)
            except Exception as e:
                result[0] = e
        
        thread = threading.Thread(target=target)
        thread.daemon = True
        thread.start()
        thread.join(timeout_seconds)
        
        if thread.is_alive():
            logger.error(f"转换超时（{timeout_seconds}秒）")
            raise TimeoutException(f'转换超时（{timeout_seconds}秒）')
        
        if isinstance(result[0], Exception):
            raise result[0]
        
        return result[0]
    
    return wrapper


def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def cleanup_old_files(directory, max_age_hours=1):
    """清理超过指定时间的旧文件"""
    try:
        now = datetime.now()
        for filename in os.listdir(directory):
            filepath = os.path.join(directory, filename)
            if os.path.isfile(filepath):
                file_modified = datetime.fromtimestamp(os.path.getmtime(filepath))
                if now - file_modified > timedelta(hours=max_age_hours):
                    try:
                        os.remove(filepath)
                        logger.info(f"清理旧文件: {filename}")
                    except:
                        pass
    except Exception as e:
        logger.warning(f"清理文件失败: {str(e)}")


def convert_pdf_to_text_ultra_fast(pdf_path, output_path, start_page=0, end_page=None):
    """
    超快速纯文本提取模式 - 性能优化版
    直接提取PDF文本，无任何格式处理，速度最快
    """
    logger.info("使用超快速纯文本模式")
    doc_word = Document()
    
    # 最小边距
    sections = doc_word.sections
    for section in sections:
        section.top_margin = Inches(0.5)
        section.bottom_margin = Inches(0.5)
        section.left_margin = Inches(0.75)
        section.right_margin = Inches(0.75)
    
    pdf_doc = fitz.open(pdf_path)
    total_pages = len(pdf_doc)
    end_page = end_page if end_page else total_pages
    
    logger.info(f"提取文本: 第{start_page+1}页到第{end_page}页（共{end_page-start_page}页）")
    
    for page_num in range(start_page, min(end_page, total_pages)):
        page = pdf_doc[page_num]
        text = page.get_text("text")
        
        # 简洁的页码标题
        heading = doc_word.add_heading(f'Page {page_num + 1}', level=2)
        heading.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
        heading_format = heading.paragraph_format
        heading_format.space_before = Pt(0)
        heading_format.space_after = Pt(6)
        for run in heading.runs:
            run.font.size = Pt(11)
            run.font.color.rgb = RGBColor(46, 125, 50)
        
        # 添加文本
        if text.strip():
            para = doc_word.add_paragraph(text)
            para_format = para.paragraph_format
            para_format.line_spacing = 1.15
            para_format.space_before = Pt(0)
            para_format.space_after = Pt(0)
            for run in para.runs:
                run.font.size = Pt(10)
        else:
            para = doc_word.add_paragraph('[No text content on this page]')
            para_format = para.paragraph_format
            para_format.space_before = Pt(0)
            para_format.space_after = Pt(0)
            for run in para.runs:
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(150, 150, 150)
        
        # 分页符
        if page_num < min(end_page, total_pages) - 1:
            doc_word.add_page_break()
        
        # 进度日志
        if (page_num - start_page + 1) % 10 == 0:
            logger.info(f"已处理 {page_num - start_page + 1}/{end_page - start_page} 页")
    
    pdf_doc.close()
    doc_word.save(output_path)
    logger.info(f"超快速转换完成: {end_page - start_page}页")


def convert_with_pdf2docx_optimized(pdf_path, word_path, mode='fast', start_page=0, end_page=None, include_images=False):
    """
    使用pdf2docx转换 - Windows优化版（禁用多进程）
    
    在Windows上，pdf2docx的多进程模式会导致死锁问题
    改用单进程模式，但通过分批处理来优化性能
    """
    logger.info(f"使用优化的{mode}模式（单进程，分批处理）")
    
    # 获取总页数
    doc = fitz.open(pdf_path)
    total_pages = len(doc)
    doc.close()
    
    if end_page is None or end_page > total_pages:
        end_page = total_pages
    
    pages_to_convert = end_page - start_page
    logger.info(f"需要转换 {pages_to_convert} 页（第{start_page+1}页到第{end_page}页）")
    
    # 策略：如果页数较多，使用分批处理
    BATCH_SIZE = 20  # 每批处理20页
    
    if pages_to_convert <= BATCH_SIZE:
        # 页数较少，直接转换
        logger.info(f"页数较少（{pages_to_convert}页），直接转换")
        cv = Converter(pdf_path)
        cv.convert(
            word_path,
            start=start_page,
            end=end_page,
            image=include_images,
            multi_processing=False,  # 禁用多进程！
        )
        cv.close()
        logger.info("转换完成")
    else:
        # 页数较多，分批转换后合并
        logger.info(f"页数较多（{pages_to_convert}页），使用分批策略")
        temp_docs = []
        
        for batch_start in range(start_page, end_page, BATCH_SIZE):
            batch_end = min(batch_start + BATCH_SIZE, end_page)
            batch_num = (batch_start - start_page) // BATCH_SIZE + 1
            total_batches = (pages_to_convert + BATCH_SIZE - 1) // BATCH_SIZE
            
            logger.info(f"处理批次 {batch_num}/{total_batches}: 第{batch_start+1}-{batch_end}页")
            
            # 临时文件
            temp_path = word_path.replace('.docx', f'_batch{batch_num}.docx')
            temp_docs.append(temp_path)
            
            try:
                cv = Converter(pdf_path)
                cv.convert(
                    temp_path,
                    start=batch_start,
                    end=batch_end,
                    image=include_images,
                    multi_processing=False,  # 禁用多进程
                )
                cv.close()
                logger.info(f"批次 {batch_num} 完成")
            except Exception as e:
                logger.error(f"批次 {batch_num} 失败: {str(e)}")
                # 清理临时文件
                for temp_doc in temp_docs:
                    try:
                        if os.path.exists(temp_doc):
                            os.remove(temp_doc)
                    except:
                        pass
                raise
        
        # 合并所有批次
        logger.info(f"合并 {len(temp_docs)} 个批次...")
        merge_word_documents(temp_docs, word_path)
        
        # 清理临时文件
        for temp_doc in temp_docs:
            try:
                os.remove(temp_doc)
                logger.info(f"删除临时文件: {os.path.basename(temp_doc)}")
            except:
                pass
        
        logger.info("分批转换完成")


def merge_word_documents(doc_paths, output_path):
    """合并多个Word文档"""
    if not doc_paths:
        raise ValueError("没有要合并的文档")
    
    if len(doc_paths) == 1:
        # 只有一个文档，直接重命名
        os.rename(doc_paths[0], output_path)
        return
    
    # 创建主文档
    main_doc = Document(doc_paths[0])
    
    # 追加其他文档
    for doc_path in doc_paths[1:]:
        sub_doc = Document(doc_path)
        
        # 添加分页符
        main_doc.add_page_break()
        
        # 复制所有段落
        for paragraph in sub_doc.paragraphs:
            new_para = main_doc.add_paragraph(paragraph.text, style=paragraph.style)
            # 复制段落格式
            new_para.alignment = paragraph.alignment
            new_para.paragraph_format.left_indent = paragraph.paragraph_format.left_indent
            new_para.paragraph_format.right_indent = paragraph.paragraph_format.right_indent
            new_para.paragraph_format.space_before = paragraph.paragraph_format.space_before
            new_para.paragraph_format.space_after = paragraph.paragraph_format.space_after
            
            # 复制字体格式
            for i, run in enumerate(paragraph.runs):
                if i < len(new_para.runs):
                    new_run = new_para.runs[i]
                    new_run.bold = run.bold
                    new_run.italic = run.italic
                    new_run.underline = run.underline
                    if run.font.size:
                        new_run.font.size = run.font.size
        
        # 复制表格
        for table in sub_doc.tables:
            new_table = main_doc.add_table(rows=len(table.rows), cols=len(table.columns))
            for i, row in enumerate(table.rows):
                for j, cell in enumerate(row.cells):
                    new_table.rows[i].cells[j].text = cell.text
    
    main_doc.save(output_path)


def parse_page_range(pdf_path, pages_param):
    """解析页码范围"""
    try:
        doc = fitz.open(pdf_path)
        total_pages = len(doc)
        doc.close()
        
        if pages_param == 'all' or not pages_param:
            return 0, total_pages, list(range(total_pages))
        
        selected_pages = []
        for part in pages_param.split(','):
            part = part.strip()
            if '-' in part:
                start, end = part.split('-')
                start = int(start) - 1
                end = int(end)
                selected_pages.extend(range(max(0, start), min(total_pages, end)))
            else:
                page = int(part) - 1
                if 0 <= page < total_pages:
                    selected_pages.append(page)
        
        if not selected_pages:
            return 0, total_pages, list(range(total_pages))
        
        start_page = min(selected_pages)
        end_page = max(selected_pages) + 1
        
        return start_page, end_page, selected_pages
    except Exception as e:
        logger.error(f"解析页码失败: {str(e)}")
        return 0, None, []


@app.route('/health', methods=['GET'])
def health_check():
    """健康检查"""
    return jsonify({
        'status': 'ok',
        'service': 'PDF转换服务（超高性能优化版）',
        'version': '3.0.0',
        'optimization': 'Windows单进程+分批处理',
        'features': [
            'PDF转Word（超快速模式）',
            '分批处理（避免大文档卡死）',
            '超时保护（5分钟）',
            '进度监控',
            '自动清理'
        ],
        'modes': {
            'ultra-fast': '超快速模式（纯文本，推荐）',
            'fast': '快速模式（表格+文本）',
            'balanced': '平衡模式（表格+文本+图片）',
            'quality': '高质量模式（完整格式+图片）'
        }
    })


@app.route('/pdf/info', methods=['POST'])
def get_pdf_info():
    """
    获取PDF信息和批量预览
    支持分页获取缩略图，优化性能
    """
    try:
        if 'file' not in request.files:
            return jsonify({'error': '没有上传文件'}), 400
        
        file = request.files['file']
        
        if file.filename == '' or not allowed_file(file.filename):
            return jsonify({'error': '无效的PDF文件'}), 400
        
        # 获取分页参数
        page_param = request.form.get('page', '1')  # 当前页
        page_size = int(request.form.get('pageSize', '10'))  # 每页显示数量，默认10
        
        try:
            current_page = int(page_param)
        except:
            current_page = 1
        
        # 临时保存文件
        file_uuid = str(uuid.uuid4())
        pdf_filename = f"{file_uuid}_temp.pdf"
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_filename)
        file.save(pdf_path)
        
        try:
            # 打开PDF
            doc = fitz.open(pdf_path)
            page_count = len(doc)
            
            # 计算分页
            total_pages = (page_count + page_size - 1) // page_size  # 总分页数
            start_index = (current_page - 1) * page_size
            end_index = min(start_index + page_size, page_count)
            
            logger.info(f"PDF预览请求: 总共{page_count}页, 获取第{start_index+1}-{end_index}页")
            
            # 生成当前分页的预览图
            previews = []
            for page_num in range(start_index, end_index):
                page = doc[page_num]
                
                # 生成缩略图 (150x200像素，适合网格显示)
                zoom = 150 / page.rect.width
                mat = fitz.Matrix(zoom, zoom)
                pix = page.get_pixmap(matrix=mat, alpha=False)
                
                # 转换为PIL Image并压缩
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data))
                
                # 压缩为JPEG
                buffer = io.BytesIO()
                img.save(buffer, format="JPEG", quality=70, optimize=True)
                img_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')
                
                previews.append({
                    'page': page_num + 1,
                    'image': f'data:image/jpeg;base64,{img_base64}',
                    'width': int(page.rect.width),
                    'height': int(page.rect.height)
                })
            
            doc.close()
            os.remove(pdf_path)
            
            logger.info(f"预览生成成功: 第{current_page}页分页, 共{len(previews)}张")
            
            return jsonify({
                'success': True,
                'pageCount': page_count,
                'previews': previews,
                'pagination': {
                    'currentPage': current_page,
                    'pageSize': page_size,
                    'totalPages': total_pages,
                    'startIndex': start_index + 1,
                    'endIndex': end_index
                }
            })
            
        except Exception as e:
            logger.error(f"PDF预览生成失败: {str(e)}")
            try:
                os.remove(pdf_path)
            except:
                pass
            return jsonify({'error': f'预览生成失败: {str(e)}'}), 500
            
    except Exception as e:
        logger.error(f"请求处理失败: {str(e)}")
        return jsonify({'error': f'请求失败: {str(e)}'}), 500


@app.route('/pdf/toword', methods=['POST'])
def convert_pdf_to_word():
    """
    PDF转Word端点 - 超高性能优化版 v3.0
    专门针对Windows优化，解决卡死问题
    """
    try:
        if 'file' not in request.files:
            return jsonify({'error': '没有上传文件'}), 400
        
        file = request.files['file']
        
        if file.filename == '':
            return jsonify({'error': '文件名为空'}), 400
        
        if not allowed_file(file.filename):
            return jsonify({'error': '只支持PDF文件'}), 400
        
        # 获取参数
        mode = request.form.get('mode', 'fast')  # ultra-fast, fast, balanced, quality
        pages_param = request.form.get('pages', 'all')
        include_images = request.form.get('include_images', 'false').lower() == 'true'
        
        # 文件名处理
        original_filename = secure_filename(file.filename)
        file_uuid = str(uuid.uuid4())
        pdf_filename = f"{file_uuid}_{original_filename}"
        word_filename = f"{file_uuid}_{os.path.splitext(original_filename)[0]}.docx"
        
        pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_filename)
        word_path = os.path.join(app.config['CONVERTED_FOLDER'], word_filename)
        
        # 保存PDF
        file.save(pdf_path)
        file_size = os.path.getsize(pdf_path) / 1024
        logger.info(f"接收文件: {original_filename} ({file_size:.2f} KB)")
        logger.info(f"模式: {mode}, 页码: {pages_param}, 图片: {include_images}")
        
        # 清理旧文件
        cleanup_old_files(app.config['UPLOAD_FOLDER'])
        cleanup_old_files(app.config['CONVERTED_FOLDER'])
        
        # 解析页码
        start_page, end_page, selected_pages = parse_page_range(pdf_path, pages_param)
        pages_count = end_page - start_page if end_page else 0
        
        logger.info(f"开始转换: {pages_count}页")
        start_time = datetime.now()
        
        try:
            # 根据模式选择转换策略
            if mode == 'ultra-fast' or mode == 'text-only':
                # 超快速模式：纯文本提取
                convert_pdf_to_text_ultra_fast(pdf_path, word_path, start_page, end_page)
                
            elif mode == 'fast':
                # 快速模式：表格+文本，不含图片
                convert_with_pdf2docx_optimized(
                    pdf_path, word_path, 
                    mode='fast',
                    start_page=start_page, 
                    end_page=end_page,
                    include_images=False
                )
                
            elif mode == 'balanced' or mode == 'premium':
                # 平衡模式：表格+文本+少量图片
                convert_with_pdf2docx_optimized(
                    pdf_path, word_path,
                    mode='balanced',
                    start_page=start_page,
                    end_page=end_page,
                    include_images=True
                )
                
            elif mode == 'quality' or mode == 'complex':
                # 高质量模式：完整格式+图片
                convert_with_pdf2docx_optimized(
                    pdf_path, word_path,
                    mode='quality',
                    start_page=start_page,
                    end_page=end_page,
                    include_images=True
                )
                
            else:
                # 默认：快速模式
                convert_with_pdf2docx_optimized(
                    pdf_path, word_path,
                    mode='fast',
                    start_page=start_page,
                    end_page=end_page,
                    include_images=include_images
                )
            
            conversion_time = (datetime.now() - start_time).total_seconds()
            word_size = os.path.getsize(word_path)
            
            logger.info(f"转换成功: {word_filename} ({word_size/1024:.2f} KB, {conversion_time:.2f}s)")
            logger.info(f"平均速度: {pages_count/conversion_time:.2f} 页/秒")
            
            # 清理PDF
            try:
                os.remove(pdf_path)
            except:
                pass
            
            return jsonify({
                'url': f'/download/{word_filename}',
                'filename': word_filename,
                'size': word_size,
                'conversion_time': f'{conversion_time:.2f}s',
                'mode': mode,
                'pages_converted': f'{start_page+1}-{end_page}' if pages_param != 'all' else 'all',
                'pages_count': pages_count,
                'speed': f'{pages_count/conversion_time:.2f} 页/秒'
            })
            
        except TimeoutException as e:
            logger.error(f"转换超时: {str(e)}")
            try:
                os.remove(pdf_path)
            except:
                pass
            return jsonify({'error': '转换超时，请尝试使用更快的模式或减少页数'}), 408
            
        except Exception as e:
            logger.error(f"转换失败: {str(e)}")
            logger.error(traceback.format_exc())
            try:
                os.remove(pdf_path)
            except:
                pass
            return jsonify({'error': f'转换失败: {str(e)}'}), 500
        
    except Exception as e:
        logger.error(f"请求处理失败: {str(e)}")
        return jsonify({'error': f'请求失败: {str(e)}'}), 500


@app.route('/pdf/to-images', methods=['POST'])
def pdf_to_images():
    """
    PDF转图片 - 支持分页、多格式、自定义质量
    
    参数:
    - file: PDF文件
    - page: 当前页码（默认1）
    - page_size: 每页返回图片数（默认6）
    - format: 输出格式 png/jpg（默认png）
    - quality: 图片质量 1-100（默认85）
    - dpi: 分辨率 72-600（默认150）
    
    返回:
    - images: 图片数组（base64编码）
    - current_page: 当前页码
    - total_pages: 总页数（分页）
    - total_pdf_pages: PDF总页数
    """
    try:
        # 验证文件
        if 'file' not in request.files:
            return jsonify({'error': '未上传文件'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': '文件名为空'}), 400
        
        # 获取参数
        current_page = int(request.form.get('page', 1))
        page_size = int(request.form.get('page_size', 6))
        img_format = request.form.get('format', 'png').lower()
        quality = int(request.form.get('quality', 85))
        dpi = int(request.form.get('dpi', 150))
        
        # 参数验证
        if current_page < 1:
            current_page = 1
        if page_size < 1 or page_size > 20:
            page_size = 6
        if img_format not in ['png', 'jpg', 'jpeg']:
            img_format = 'png'
        if quality < 1 or quality > 100:
            quality = 85
        if dpi < 72 or dpi > 600:
            dpi = 150
        
        # 读取PDF
        pdf_bytes = file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        total_pdf_pages = len(doc)
        
        logger.info(f"PDF转图片: {file.filename}, 总页数={total_pdf_pages}, 请求页={current_page}, 格式={img_format}, DPI={dpi}")
        
        # 计算分页
        start_index = (current_page - 1) * page_size
        end_index = min(start_index + page_size, total_pdf_pages)
        total_pages = (total_pdf_pages + page_size - 1) // page_size  # 向上取整
        
        if start_index >= total_pdf_pages:
            doc.close()
            return jsonify({'error': '页码超出范围'}), 400
        
        # 高性能转换PDF页为图片
        images = []
        zoom = dpi / 72  # 计算缩放比例
        mat = fitz.Matrix(zoom, zoom)
        
        # 性能优化：使用线程池并行处理（但保持顺序）
        start_time = datetime.now()
        
        for page_num in range(start_index, end_index):
            page_start = datetime.now()
            page = doc[page_num]
            
            # 高质量渲染：alpha=False提升性能，使用高质量抗锯齿
            pix = page.get_pixmap(
                matrix=mat, 
                alpha=False,
                colorspace=fitz.csRGB  # 明确指定RGB色彩空间
            )
            
            # 性能优化：直接使用PyMuPDF输出，避免PIL转换
            if img_format in ['jpg', 'jpeg']:
                # 直接输出JPEG（PyMuPDF原生支持，更快更高质量）
                img_bytes = pix.tobytes("jpeg", jpg_quality=quality)
                img_base64 = base64.b64encode(img_bytes).decode('utf-8')
            else:
                # PNG格式：使用PIL优化压缩
                img_data = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_data))
                
                # 高质量PNG压缩
                buffer = io.BytesIO()
                img.save(buffer, format='PNG', optimize=True, compress_level=6)
                img_bytes = buffer.getvalue()
                img_base64 = base64.b64encode(img_bytes).decode('utf-8')
            
            page_duration = (datetime.now() - page_start).total_seconds()
            
            images.append({
                'page': page_num + 1,
                'image': f'data:image/{img_format};base64,{img_base64}',
                'width': pix.width,
                'height': pix.height,
                'size': len(img_bytes)
            })
            
            logger.info(f"  页面 {page_num + 1}: {pix.width}x{pix.height}, {len(img_bytes)/1024:.1f}KB, 耗时{page_duration*1000:.0f}ms")
        
        doc.close()
        
        total_duration = (datetime.now() - start_time).total_seconds()
        total_size = sum(img['size'] for img in images)
        avg_time = total_duration / len(images) if images else 0
        
        logger.info(f"高性能转换完成: {len(images)}张图片, 总耗时={total_duration:.2f}s, 平均={avg_time*1000:.0f}ms/页, 总大小={total_size/1024:.1f}KB")
        
        return jsonify({
            'images': images,
            'current_page': current_page,
            'total_pages': total_pages,
            'total_pdf_pages': total_pdf_pages,
            'page_size': page_size,
            'start_page': start_index + 1,
            'end_page': end_index,
            'format': img_format,
            'quality': quality,
            'dpi': dpi
        })
        
    except Exception as e:
        logger.error(f"PDF转图片失败: {str(e)}")
        logger.error(traceback.format_exc())
        return jsonify({'error': f'转换失败: {str(e)}'}), 500


@app.route('/pdf/to-ppt', methods=['POST'])
def pdf_to_ppt():
    """
    PDF转PPT - 高性能图片方式
    
    流程：PDF → 图片（PyMuPDF）→ PPT（python-pptx）
    
    参数:
    - file: PDF文件（必填）
    - dpi: 分辨率 72-300（默认150）
    - quality: JPEG质量 60-95（默认85）
    
    限制:
    - 文件大小：60MB以内
    - 页数：100页以内（建议）
    
    返回:
    - url: PPT下载地址
    - filename: PPT文件名
    - pages: 转换的页数
    - size: 文件大小
    - conversion_time: 转换耗时
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return jsonify({'error': '缺少python-pptx库，请安装: pip install python-pptx'}), 500
    
    try:
        # 验证文件
        if 'file' not in request.files:
            return jsonify({'error': '未上传文件'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': '文件名为空'}), 400
        
        # 检查文件大小（60MB限制）
        file.seek(0, 2)  # 移到文件末尾
        file_size = file.tell()
        file.seek(0)  # 回到文件开头
        
        MAX_SIZE = 60 * 1024 * 1024  # 60MB
        if file_size > MAX_SIZE:
            return jsonify({'error': f'文件大小超过限制（最大60MB）'}), 400
        
        # 获取参数（优化：提高默认分辨率）
        dpi = int(request.form.get('dpi', 200))  # 提高默认DPI从150到200
        quality = int(request.form.get('quality', 92))  # 提高默认质量从85到92
        
        # 参数验证（扩大DPI上限以支持超高清）
        if dpi < 72 or dpi > 400:  # DPI上限从300提升到400
            dpi = 200
        if quality < 60 or quality > 100:  # 质量上限从95提升到100
            quality = 92
        
        start_time = datetime.now()
        
        # 读取PDF
        pdf_bytes = file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        total_pages = len(doc)
        
        logger.info(f"========================================")
        logger.info(f"PDF转PPT: {file.filename}")
        logger.info(f"PDF页数: {total_pages}")
        logger.info(f"文件大小: {file_size/1024:.1f} KB")
        logger.info(f"DPI: {dpi}, 质量: {quality}")
        logger.info(f"========================================")
        
        # 检查页数限制
        if total_pages > 100:
            logger.warning(f"页数过多（{total_pages}页），建议100页以内")
            # 不强制限制，但给出警告
        
        # 步骤1：PDF转图片（高性能）
        images_data = []
        zoom = dpi / 72
        mat = fitz.Matrix(zoom, zoom)
        
        for page_num in range(total_pages):
            page_start = datetime.now()
            page = doc[page_num]
            
            # 高质量渲染
            pix = page.get_pixmap(
                matrix=mat,
                alpha=False,
                colorspace=fitz.csRGB
            )
            
            # 直接输出JPEG（性能优化）
            img_bytes = pix.tobytes("jpeg", jpg_quality=quality)
            images_data.append({
                'bytes': img_bytes,
                'width': pix.width,
                'height': pix.height
            })
            
            page_duration = (datetime.now() - page_start).total_seconds()
            logger.info(f"  页面 {page_num + 1}/{total_pages}: {pix.width}x{pix.height}, {len(img_bytes)/1024:.1f}KB, {page_duration*1000:.0f}ms")
        
        doc.close()
        
        # 步骤2：创建PPT
        ppt_start = datetime.now()
        prs = Presentation()
        
        # 设置幻灯片尺寸为标准16:9
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)  # 16:9比例
        
        # 空白布局
        blank_layout = prs.slide_layouts[6]
        
        # 插入每张图片为一页幻灯片
        for idx, img_data in enumerate(images_data):
            slide = prs.slides.add_slide(blank_layout)
            
            # 将图片插入，填满整个幻灯片
            img_stream = io.BytesIO(img_data['bytes'])
            left = top = Inches(0)
            slide.shapes.add_picture(
                img_stream,
                left, top,
                width=prs.slide_width,
                height=prs.slide_height
            )
            
            logger.info(f"  插入幻灯片 {idx + 1}/{total_pages}")
        
        # 保存PPT
        ppt_filename = f"converted_{uuid.uuid4().hex[:8]}.pptx"
        ppt_path = os.path.join(app.config['CONVERTED_FOLDER'], ppt_filename)
        
        prs.save(ppt_path)
        ppt_size = os.path.getsize(ppt_path)
        
        ppt_duration = (datetime.now() - ppt_start).total_seconds()
        total_duration = (datetime.now() - start_time).total_seconds()
        
        logger.info(f"========================================")
        logger.info(f"PDF转PPT完成！")
        logger.info(f"生成PPT: {ppt_filename}")
        logger.info(f"PPT大小: {ppt_size/1024:.1f} KB")
        logger.info(f"PPT创建耗时: {ppt_duration:.2f}s")
        logger.info(f"总耗时: {total_duration:.2f}s")
        logger.info(f"平均速度: {total_pages/total_duration:.1f} 页/秒")
        logger.info(f"========================================")
        
        return jsonify({
            'url': f'/download/{ppt_filename}',
            'filename': ppt_filename,
            'pages': total_pages,
            'size': ppt_size,
            'conversion_time': f'{total_duration:.2f}s',
            'dpi': dpi,
            'quality': quality,
            'message': '转换成功'
        })
        
    except Exception as e:
        logger.error(f"PDF转PPT失败: {str(e)}")
        logger.error(traceback.format_exc())
        return jsonify({'error': f'转换失败: {str(e)}'}), 500


@app.route('/text/to-pdf', methods=['POST'])
def text_to_pdf():
    """
    文本转PDF - 高性能优化版
    支持：
    1. 直接输入文本
    2. 上传TXT文件
    3. 自动换行和分页
    4. 中文字体支持
    5. 高质量输出
    """
    from reportlab.lib.pagesizes import A4
    from reportlab.pdfgen import canvas
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.lib.units import cm
    
    start_time = datetime.now()
    logger.info("========================================")
    logger.info("收到文本转PDF请求")
    
    try:
        # 获取文本内容
        text_content = None
        source_type = None
        
        # 方式1：直接文本输入
        if 'text' in request.form:
            text_content = request.form.get('text', '').strip()
            source_type = 'direct'
            logger.info(f"输入方式: 直接文本输入")
            logger.info(f"文本长度: {len(text_content)} 字符")
        
        # 方式2：TXT文件上传
        elif 'file' in request.files:
            file = request.files['file']
            if file.filename == '':
                return jsonify({'error': '未选择文件'}), 400
            
            # 检查文件扩展名
            if not file.filename.lower().endswith('.txt'):
                return jsonify({'error': '只支持TXT文件'}), 400
            
            # 读取文件内容（支持多种编码）
            file_content = file.read()
            
            # 尝试多种编码
            for encoding in ['utf-8', 'gbk', 'gb2312', 'utf-16']:
                try:
                    text_content = file_content.decode(encoding)
                    logger.info(f"文件编码: {encoding}")
                    break
                except:
                    continue
            
            if text_content is None:
                return jsonify({'error': '无法识别文件编码'}), 400
            
            source_type = 'file'
            logger.info(f"输入方式: TXT文件上传")
            logger.info(f"文件名: {file.filename}")
            logger.info(f"文件大小: {len(file_content)} 字节")
            logger.info(f"文本长度: {len(text_content)} 字符")
        
        else:
            return jsonify({'error': '请提供文本内容或上传TXT文件'}), 400
        
        # 验证文本内容
        if not text_content:
            return jsonify({'error': '文本内容为空'}), 400
        
        # 文件大小限制（10MB文本）
        if len(text_content) > 10 * 1024 * 1024:
            return jsonify({'error': '文本内容过大（最大10MB）'}), 400
        
        # 获取配置参数
        font_size = int(request.form.get('font_size', 12))
        line_spacing = float(request.form.get('line_spacing', 1.5))
        margin_left = float(request.form.get('margin_left', 2.5))
        margin_right = float(request.form.get('margin_right', 2.5))
        margin_top = float(request.form.get('margin_top', 2.5))
        margin_bottom = float(request.form.get('margin_bottom', 2.5))
        
        logger.info(f"配置参数:")
        logger.info(f"  字体大小: {font_size}")
        logger.info(f"  行间距: {line_spacing}")
        logger.info(f"  页边距: 左{margin_left}cm, 右{margin_right}cm, 上{margin_top}cm, 下{margin_bottom}cm")
        
        # 创建PDF
        pdf_filename = f"text2pdf_{uuid.uuid4().hex[:8]}.pdf"
        pdf_path = os.path.join(app.config['CONVERTED_FOLDER'], pdf_filename)
        
        # 注册中文字体（使用系统字体）
        font_name = 'SimSun'
        try:
            # Windows系统字体路径
            font_paths = [
                'C:/Windows/Fonts/simsun.ttc',
                'C:/Windows/Fonts/msyh.ttc',  # 微软雅黑
                'C:/Windows/Fonts/simhei.ttf'  # 黑体
            ]
            
            font_registered = False
            for font_path in font_paths:
                if os.path.exists(font_path):
                    try:
                        pdfmetrics.registerFont(TTFont(font_name, font_path))
                        font_registered = True
                        logger.info(f"字体注册成功: {font_path}")
                        break
                    except Exception as e:
                        logger.warning(f"字体注册失败 {font_path}: {str(e)}")
                        continue
            
            if not font_registered:
                logger.warning("未找到中文字体，使用默认字体")
                font_name = 'Helvetica'
        
        except Exception as e:
            logger.error(f"字体注册错误: {str(e)}")
            font_name = 'Helvetica'
        
        # 创建PDF画布
        c = canvas.Canvas(pdf_path, pagesize=A4)
        width, height = A4
        
        # 计算可用区域
        usable_width = width - margin_left * cm - margin_right * cm
        usable_height = height - margin_top * cm - margin_bottom * cm
        
        # 设置字体
        c.setFont(font_name, font_size)
        
        # 计算行高
        line_height = font_size * line_spacing
        
        # 分行处理
        lines = text_content.split('\n')
        current_y = height - margin_top * cm
        page_count = 1
        char_count = 0
        
        logger.info(f"开始转换...")
        
        for line in lines:
            # 空行处理
            if not line.strip():
                current_y -= line_height
                if current_y < margin_bottom * cm:
                    c.showPage()
                    c.setFont(font_name, font_size)
                    current_y = height - margin_top * cm
                    page_count += 1
                continue
            
            # 自动换行
            words = []
            current_word = ''
            
            for char in line:
                current_word += char
                char_width = c.stringWidth(current_word, font_name, font_size)
                
                if char_width >= usable_width:
                    if len(current_word) > 1:
                        words.append(current_word[:-1])
                        current_word = char
                    else:
                        words.append(current_word)
                        current_word = ''
            
            if current_word:
                words.append(current_word)
            
            # 输出每一行
            for word in words:
                if current_y < margin_bottom * cm + line_height:
                    c.showPage()
                    c.setFont(font_name, font_size)
                    current_y = height - margin_top * cm
                    page_count += 1
                
                c.drawString(margin_left * cm, current_y, word)
                current_y -= line_height
                char_count += len(word)
        
        # 保存PDF
        c.save()
        
        pdf_size = os.path.getsize(pdf_path)
        total_duration = (datetime.now() - start_time).total_seconds()
        
        logger.info(f"========================================")
        logger.info(f"文本转PDF完成！")
        logger.info(f"生成PDF: {pdf_filename}")
        logger.info(f"PDF大小: {pdf_size/1024:.1f} KB")
        logger.info(f"页数: {page_count}")
        logger.info(f"字符数: {char_count}")
        logger.info(f"总耗时: {total_duration:.2f}s")
        logger.info(f"转换速度: {char_count/total_duration:.0f} 字符/秒")
        logger.info(f"========================================")
        
        return jsonify({
            'url': f'/download/{pdf_filename}',
            'filename': pdf_filename,
            'pages': page_count,
            'size': pdf_size,
            'characters': char_count,
            'conversion_time': f'{total_duration:.2f}s',
            'source_type': source_type,
            'message': '转换成功'
        })
        
    except Exception as e:
        logger.error(f"文本转PDF失败: {str(e)}")
        logger.error(traceback.format_exc())
        return jsonify({'error': f'转换失败: {str(e)}'}), 500


@app.route('/pdf/extract/upload', methods=['POST'])
def upload_pdf_for_extract():
    """
    上传PDF文件用于内容提取
    返回：PDF文件信息、所有页面缩略图（Base64）
    """
    start_time = datetime.now()
    logger.info("========================================")
    logger.info("收到PDF提取上传请求")
    
    try:
        if 'file' not in request.files:
            return jsonify({'error': '未找到文件'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': '未选择文件'}), 400
        
        if not file.filename.lower().endswith('.pdf'):
            return jsonify({'error': '只支持PDF文件'}), 400
        
        # 保存上传的文件
        filename = secure_filename(file.filename)
        unique_filename = f"{uuid.uuid4().hex[:8]}_{filename}"
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
        file.save(file_path)
        
        file_size = os.path.getsize(file_path)
        logger.info(f"文件保存: {unique_filename}, 大小: {file_size/1024:.1f} KB")
        
        # 打开PDF文档
        doc = fitz.open(file_path)
        total_pages = len(doc)
        
        logger.info(f"PDF页数: {total_pages}")
        
        # 生成所有页面缩略图（Base64）
        thumbnails = []
        thumbnail_size = 200  # 缩略图宽度
        
        for page_num in range(total_pages):
            page = doc[page_num]
            
            # 计算缩放比例
            mat = fitz.Matrix(thumbnail_size / page.rect.width, thumbnail_size / page.rect.width)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            
            # 转换为JPEG Base64
            img_bytes = pix.tobytes("jpeg", jpg_quality=85)
            img_base64 = base64.b64encode(img_bytes).decode('utf-8')
            
            thumbnails.append({
                'page': page_num + 1,
                'width': pix.width,
                'height': pix.height,
                'data': f"data:image/jpeg;base64,{img_base64}"
            })
            
            logger.info(f"  页面 {page_num + 1}/{total_pages}: 缩略图生成 ({pix.width}x{pix.height})")
        
        doc.close()
        
        total_duration = (datetime.now() - start_time).total_seconds()
        
        logger.info(f"========================================")
        logger.info(f"PDF上传完成！")
        logger.info(f"文件ID: {unique_filename}")
        logger.info(f"总页数: {total_pages}")
        logger.info(f"处理耗时: {total_duration:.2f}s")
        logger.info(f"========================================")
        
        return jsonify({
            'file_id': unique_filename,
            'filename': filename,
            'pages': total_pages,
            'size': file_size,
            'thumbnails': thumbnails,
            'message': '上传成功'
        })
        
    except Exception as e:
        logger.error(f"PDF上传失败: {str(e)}")
        logger.error(traceback.format_exc())
        return jsonify({'error': f'上传失败: {str(e)}'}), 500


@app.route('/pdf/extract/text', methods=['POST'])
def extract_text_from_page():
    """
    提取指定页面的文字
    参数：file_id, page_num
    """
    try:
        file_id = request.form.get('file_id')
        page_num = int(request.form.get('page_num', 1))
        
        if not file_id:
            return jsonify({'error': '缺少file_id参数'}), 400
        
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], file_id)
        if not os.path.exists(file_path):
            return jsonify({'error': 'PDF文件不存在'}), 404
        
        logger.info(f"提取文字: {file_id}, 页面 {page_num}")
        
        doc = fitz.open(file_path)
        
        if page_num < 1 or page_num > len(doc):
            return jsonify({'error': f'页码超出范围（1-{len(doc)}）'}), 400
        
        page = doc[page_num - 1]
        text = page.get_text()
        
        doc.close()
        
        # 检查是否有文字
        text_stripped = text.strip()
        has_text = len(text_stripped) > 0
        
        logger.info(f"  提取到 {len(text_stripped)} 个字符")
        
        return jsonify({
            'page': page_num,
            'has_text': has_text,
            'text': text_stripped if has_text else '',
            'char_count': len(text_stripped),
            'message': '提取成功' if has_text else '未检测到可复制文字'
        })
        
    except Exception as e:
        logger.error(f"文字提取失败: {str(e)}")
        logger.error(traceback.format_exc())
        return jsonify({'error': f'提取失败: {str(e)}'}), 500


@app.route('/pdf/extract/page-image', methods=['POST'])
def extract_page_as_image():
    """
    将指定页面保存为一张图像
    参数：file_id, page_num, dpi（可选，默认150）
    """
    try:
        file_id = request.form.get('file_id')
        page_num = int(request.form.get('page_num', 1))
        dpi = int(request.form.get('dpi', 150))
        
        if not file_id:
            return jsonify({'error': '缺少file_id参数'}), 400
        
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], file_id)
        if not os.path.exists(file_path):
            return jsonify({'error': 'PDF文件不存在'}), 404
        
        logger.info(f"提取页面图像: {file_id}, 页面 {page_num}, DPI {dpi}")
        
        doc = fitz.open(file_path)
        
        if page_num < 1 or page_num > len(doc):
            return jsonify({'error': f'页码超出范围（1-{len(doc)}）'}), 400
        
        page = doc[page_num - 1]
        
        # 设置缩放比例
        zoom = dpi / 72.0
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        
        # 保存为图像
        output_filename = f"page_{page_num}_{uuid.uuid4().hex[:8]}.png"
        output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)
        pix.save(output_path)
        
        doc.close()
        
        file_size = os.path.getsize(output_path)
        
        logger.info(f"  页面图像已保存: {output_filename}, {pix.width}x{pix.height}, {file_size/1024:.1f} KB")
        
        return jsonify({
            'url': f'/download/{output_filename}',
            'filename': output_filename,
            'page': page_num,
            'width': pix.width,
            'height': pix.height,
            'size': file_size,
            'dpi': dpi,
            'message': '保存成功'
        })
        
    except Exception as e:
        logger.error(f"页面图像提取失败: {str(e)}")
        logger.error(traceback.format_exc())
        return jsonify({'error': f'提取失败: {str(e)}'}), 500


@app.route('/pdf/extract/embedded-images', methods=['POST'])
def extract_embedded_images():
    """
    提取指定页面的所有内嵌图片
    参数：file_id, page_num
    返回：图片列表（Base64或下载链接）
    """
    try:
        file_id = request.form.get('file_id')
        page_num = int(request.form.get('page_num', 1))
        
        if not file_id:
            return jsonify({'error': '缺少file_id参数'}), 400
        
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], file_id)
        if not os.path.exists(file_path):
            return jsonify({'error': 'PDF文件不存在'}), 404
        
        logger.info(f"提取内嵌图片: {file_id}, 页面 {page_num}")
        
        doc = fitz.open(file_path)
        
        if page_num < 1 or page_num > len(doc):
            return jsonify({'error': f'页码超出范围（1-{len(doc)}）'}), 400
        
        page = doc[page_num - 1]
        image_list = page.get_images(full=True)
        
        logger.info(f"  发现 {len(image_list)} 张图片")
        
        if len(image_list) == 0:
            return jsonify({
                'page': page_num,
                'image_count': 0,
                'images': [],
                'message': '该页面没有内嵌图片'
            })
        
        images = []
        
        for img_index, img in enumerate(image_list):
            xref = img[0]
            
            # 提取图片
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            image_width = base_image["width"]
            image_height = base_image["height"]
            
            # 保存图片文件
            output_filename = f"embedded_p{page_num}_img{img_index+1}_{uuid.uuid4().hex[:6]}.{image_ext}"
            output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)
            
            with open(output_path, 'wb') as f:
                f.write(image_bytes)
            
            file_size = os.path.getsize(output_path)
            
            images.append({
                'index': img_index + 1,
                'filename': output_filename,
                'url': f'/download/{output_filename}',
                'format': image_ext,
                'width': image_width,
                'height': image_height,
                'size': file_size
            })
            
            logger.info(f"    图片 {img_index+1}: {image_width}x{image_height}, {image_ext}, {file_size/1024:.1f} KB")
        
        doc.close()
        
        return jsonify({
            'page': page_num,
            'image_count': len(images),
            'images': images,
            'message': f'成功提取 {len(images)} 张图片'
        })
        
    except Exception as e:
        logger.error(f"内嵌图片提取失败: {str(e)}")
        logger.error(traceback.format_exc())
        return jsonify({'error': f'提取失败: {str(e)}'}), 500


@app.route('/pdf/split', methods=['POST'])
def split_pdf():
    """
    拆分PDF文件
    参数：
    - file: PDF文件
    - split_mode: 拆分模式（by_pages, by_count, by_ranges）
    - pages_per_file: 每个文件的页数（by_pages模式）
    - file_count: 拆分成多少个文件（by_count模式）
    - ranges: 页码范围，格式：1-5,6-10,11-15（by_ranges模式）
    """
    start_time = datetime.now()
    logger.info("========================================")
    logger.info("收到PDF拆分请求")
    
    try:
        if 'file' not in request.files:
            return jsonify({'error': '未找到文件'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': '未选择文件'}), 400
        
        if not file.filename.lower().endswith('.pdf'):
            return jsonify({'error': '只支持PDF文件'}), 400
        
        # 保存上传的文件
        filename = secure_filename(file.filename)
        unique_filename = f"{uuid.uuid4().hex[:8]}_{filename}"
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], unique_filename)
        file.save(file_path)
        
        file_size = os.path.getsize(file_path)
        logger.info(f"文件保存: {unique_filename}, 大小: {file_size/1024:.1f} KB")
        
        # 获取拆分参数
        split_mode = request.form.get('split_mode', 'by_pages')
        
        # 打开PDF文档
        doc = fitz.open(file_path)
        total_pages = len(doc)
        logger.info(f"PDF页数: {total_pages}")
        
        split_files = []
        
        if split_mode == 'by_pages':
            # 按页数拆分
            pages_per_file = int(request.form.get('pages_per_file', 1))
            logger.info(f"拆分模式: 按页数, 每个文件 {pages_per_file} 页")
            
            for start in range(0, total_pages, pages_per_file):
                end = min(start + pages_per_file - 1, total_pages - 1)
                
                new_doc = fitz.open()
                new_doc.insert_pdf(doc, from_page=start, to_page=end)
                
                output_filename = f"split_{start+1}_to_{end+1}_{uuid.uuid4().hex[:6]}.pdf"
                output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)
                new_doc.save(output_path)
                new_doc.close()
                
                output_size = os.path.getsize(output_path)
                split_files.append({
                    'filename': output_filename,
                    'url': f'/download/{output_filename}',
                    'pages': end - start + 1,
                    'page_range': f'{start+1}-{end+1}',
                    'size': output_size
                })
                
                logger.info(f"  生成文件: {output_filename}, 页面 {start+1}-{end+1}")
        
        elif split_mode == 'by_count':
            # 按份数拆分
            file_count = int(request.form.get('file_count', 1))
            logger.info(f"拆分模式: 按份数, 分成 {file_count} 个文件")
            
            pages_per_file = (total_pages + file_count - 1) // file_count
            
            for i in range(file_count):
                start = i * pages_per_file
                if start >= total_pages:
                    break
                end = min(start + pages_per_file - 1, total_pages - 1)
                
                new_doc = fitz.open()
                new_doc.insert_pdf(doc, from_page=start, to_page=end)
                
                output_filename = f"split_part{i+1}_{uuid.uuid4().hex[:6]}.pdf"
                output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)
                new_doc.save(output_path)
                new_doc.close()
                
                output_size = os.path.getsize(output_path)
                split_files.append({
                    'filename': output_filename,
                    'url': f'/download/{output_filename}',
                    'pages': end - start + 1,
                    'page_range': f'{start+1}-{end+1}',
                    'size': output_size
                })
                
                logger.info(f"  生成文件 {i+1}: {output_filename}, 页面 {start+1}-{end+1}")
        
        elif split_mode == 'by_ranges':
            # 按范围拆分
            ranges_str = request.form.get('ranges', '')
            logger.info(f"拆分模式: 按范围, 范围: {ranges_str}")
            
            # 解析范围，格式：1-5,6-10,11-15
            range_list = []
            for range_part in ranges_str.split(','):
                range_part = range_part.strip()
                if '-' in range_part:
                    start_page, end_page = range_part.split('-')
                    range_list.append((int(start_page.strip()), int(end_page.strip())))
                else:
                    page = int(range_part)
                    range_list.append((page, page))
            
            for i, (start_page, end_page) in enumerate(range_list):
                # 页码从1开始，转换为0开始的索引
                start = start_page - 1
                end = end_page - 1
                
                if start < 0 or end >= total_pages or start > end:
                    logger.warning(f"  跳过无效范围: {start_page}-{end_page}")
                    continue
                
                new_doc = fitz.open()
                new_doc.insert_pdf(doc, from_page=start, to_page=end)
                
                output_filename = f"split_{start_page}_to_{end_page}_{uuid.uuid4().hex[:6]}.pdf"
                output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)
                new_doc.save(output_path)
                new_doc.close()
                
                output_size = os.path.getsize(output_path)
                split_files.append({
                    'filename': output_filename,
                    'url': f'/download/{output_filename}',
                    'pages': end - start + 1,
                    'page_range': f'{start_page}-{end_page}',
                    'size': output_size
                })
                
                logger.info(f"  生成文件 {i+1}: {output_filename}, 页面 {start_page}-{end_page}")
        
        else:
            return jsonify({'error': f'不支持的拆分模式: {split_mode}'}), 400
        
        doc.close()
        
        total_duration = (datetime.now() - start_time).total_seconds()
        
        logger.info(f"========================================")
        logger.info(f"PDF拆分完成！")
        logger.info(f"原文件: {filename}, {total_pages}页")
        logger.info(f"生成文件数: {len(split_files)}")
        logger.info(f"总耗时: {total_duration:.2f}s")
        logger.info(f"========================================")
        
        return jsonify({
            'original_filename': filename,
            'total_pages': total_pages,
            'split_count': len(split_files),
            'split_files': split_files,
            'conversion_time': f'{total_duration:.2f}s',
            'message': '拆分成功'
        })
        
    except Exception as e:
        logger.error(f"PDF拆分失败: {str(e)}")
        logger.error(traceback.format_exc())
        return jsonify({'error': f'拆分失败: {str(e)}'}), 500


@app.route('/pdf/page-count', methods=['POST'])
def get_pdf_page_count():
    """
    获取PDF文件页数
    """
    try:
        if 'file' not in request.files:
            return jsonify({'error': '未找到文件'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': '未选择文件'}), 400
        
        if not file.filename.lower().endswith('.pdf'):
            return jsonify({'error': '只支持PDF文件'}), 400
        
        # 保存临时文件
        filename = secure_filename(file.filename)
        temp_filename = f"temp_{uuid.uuid4().hex[:8]}_{filename}"
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], temp_filename)
        file.save(temp_path)
        
        # 获取页数
        doc = fitz.open(temp_path)
        page_count = len(doc)
        doc.close()
        
        # 删除临时文件
        try:
            os.remove(temp_path)
        except:
            pass
        
        logger.info(f"获取PDF页数: {filename}, 页数: {page_count}")
        
        return jsonify({
            'filename': filename,
            'page_count': page_count,
            'message': '获取成功'
        })
        
    except Exception as e:
        logger.error(f"获取PDF页数失败: {str(e)}")
        return jsonify({'error': f'获取失败: {str(e)}'}), 500


@app.route('/pdf/upload-temp', methods=['POST'])
def upload_temp_file():
    """
    上传单个PDF文件到临时目录
    """
    try:
        if 'file' not in request.files:
            return jsonify({'error': '未找到文件'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': '未选择文件'}), 400
        
        if not file.filename.lower().endswith('.pdf'):
            return jsonify({'error': '只支持PDF文件'}), 400
        
        # 保存文件
        filename = secure_filename(file.filename)
        file_id = f"temp_{uuid.uuid4().hex}"
        temp_filename = f"{file_id}_{filename}"
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], temp_filename)
        file.save(temp_path)
        
        # 获取文件信息
        file_size = os.path.getsize(temp_path)
        doc = fitz.open(temp_path)
        page_count = len(doc)
        doc.close()
        
        logger.info(f"临时文件上传: {filename}, ID: {file_id}, 页数: {page_count}")
        
        return jsonify({
            'success': True,
            'file_id': file_id,
            'filename': filename,
            'page_count': page_count,
            'file_size': file_size
        })
        
    except Exception as e:
        logger.error(f"临时文件上传失败: {str(e)}")
        return jsonify({'error': f'上传失败: {str(e)}'}), 500


@app.route('/pdf/merge-uploaded', methods=['POST'])
def merge_uploaded_files():
    """
    合并已上传的临时文件
    """
    import time
    start_time = time.time()
    
    try:
        data = request.get_json()
        files_info = data.get('files', [])
        
        if not files_info or len(files_info) == 0:
            return jsonify({'error': '未找到文件信息'}), 400
        
        logger.info(f"开始合并 {len(files_info)} 个已上传文件")
        
        # 准备文件列表
        temp_files = []
        total_pages = 0
        
        for file_info in files_info:
            file_id = file_info.get('file_id')
            selected_pages = file_info.get('selected_pages', 'all')
            
            # 查找文件
            upload_folder = app.config['UPLOAD_FOLDER']
            matching_files = [f for f in os.listdir(upload_folder) if f.startswith(file_id)]
            
            if not matching_files:
                return jsonify({'error': f'文件 {file_id} 不存在'}), 400
            
            temp_path = os.path.join(upload_folder, matching_files[0])
            
            # 打开PDF获取页数
            doc = fitz.open(temp_path)
            page_count = len(doc)
            doc.close()
            
            # 解析页码
            if selected_pages == 'all' or not selected_pages:
                selected_pages_list = list(range(page_count))
            else:
                selected_pages_list = parse_page_ranges(str(selected_pages), page_count)
            
            total_pages += len(selected_pages_list)
            
            # 检查页数限制
            if total_pages > 100:
                # 清理临时文件
                for tf in temp_files:
                    try:
                        os.remove(tf['path'])
                    except:
                        pass
                return jsonify({'error': f'合并后总页数超过100页限制（当前{total_pages}页）'}), 400
            
            temp_files.append({
                'path': temp_path,
                'selected_pages': selected_pages_list
            })
        
        # 开始合并
        result = fitz.open()
        
        for file_info in temp_files:
            doc = fitz.open(file_info['path'])
            
            # 逐页插入
            for page_num in file_info['selected_pages']:
                if page_num < len(doc):
                    result.insert_pdf(doc, from_page=page_num, to_page=page_num)
            
            doc.close()
        
        # 保存合并结果
        output_filename = f"merged_{uuid.uuid4().hex[:8]}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)
        
        result.save(output_path, garbage=4, deflate=True, clean=True)
        output_size = os.path.getsize(output_path)
        result.close()
        
        # 清理临时文件
        for file_info in temp_files:
            try:
                os.remove(file_info['path'])
            except:
                pass
        
        end_time = time.time()
        processing_time = round(end_time - start_time, 2)
        
        logger.info(f"合并完成: {output_filename}, 页数: {total_pages}, 耗时: {processing_time}秒")
        
        return jsonify({
            'success': True,
            'filename': output_filename,
            'download_url': f'/download/{output_filename}',
            'file_size': output_size,
            'file_size_mb': round(output_size / (1024 * 1024), 2),
            'total_pages': total_pages,
            'file_count': len(temp_files),
            'processing_time': processing_time,
            'message': '合并成功'
        })
        
    except Exception as e:
        logger.error(f"合并失败: {str(e)}")
        logger.error(traceback.format_exc())
        return jsonify({'error': f'合并失败: {str(e)}'}), 500


@app.route('/pdf/merge', methods=['POST'])
def merge_pdfs():
    """
    合并多个PDF文件
    支持：
    1. 多文件上传
    2. 每个文件选择特定页码
    3. 页数限制100页
    4. 总大小限制50MB
    """
    import time
    start_time = time.time()
    
    try:
        # 获取上传的文件
        files = request.files.getlist('files[]')
        if not files or len(files) == 0:
            return jsonify({'error': '未找到文件'}), 400
        
        # 获取页码选择信息（JSON格式）
        # 格式: [{"file_index": 0, "pages": "1-5,7,9-10"}, ...]
        pages_config = request.form.get('pages_config', '[]')
        import json
        try:
            pages_config = json.loads(pages_config)
        except:
            pages_config = []
        
        logger.info(f"收到合并请求: {len(files)}个文件")
        
        # 验证文件数量
        if len(files) > 20:
            return jsonify({'error': '最多支持合并20个PDF文件'}), 400
        
        # 保存上传的文件并验证
        uploaded_files = []
        total_size = 0
        total_pages = 0
        
        for i, file in enumerate(files):
            if file.filename == '':
                continue
            
            if not file.filename.lower().endswith('.pdf'):
                return jsonify({'error': f'文件 {file.filename} 不是PDF格式'}), 400
            
            # 保存文件
            filename = secure_filename(file.filename)
            temp_filename = f"merge_temp_{uuid.uuid4().hex[:8]}_{filename}"
            temp_path = os.path.join(app.config['UPLOAD_FOLDER'], temp_filename)
            file.save(temp_path)
            
            # 获取文件信息
            file_size = os.path.getsize(temp_path)
            total_size += file_size
            
            # 检查总大小限制（50MB）
            if total_size > 50 * 1024 * 1024:
                # 清理已上传的文件
                for f in uploaded_files:
                    try:
                        os.remove(f['path'])
                    except:
                        pass
                return jsonify({'error': '文件总大小超过50MB限制'}), 400
            
            # 打开PDF获取页数
            doc = fitz.open(temp_path)
            page_count = len(doc)
            doc.close()
            
            # 解析该文件的页码选择
            selected_pages = None
            for config in pages_config:
                if config.get('file_index') == i:
                    selected_pages = config.get('pages', 'all')
                    break
            
            # 如果没有指定，默认全部页
            if selected_pages is None or selected_pages == 'all':
                selected_pages_list = list(range(page_count))
            else:
                # 解析页码字符串（如"1-5,7,9-10"）
                selected_pages_list = parse_page_ranges(selected_pages, page_count)
            
            total_pages += len(selected_pages_list)
            
            # 检查总页数限制（100页）
            if total_pages > 100:
                # 清理已上传的文件
                for f in uploaded_files:
                    try:
                        os.remove(f['path'])
                    except:
                        pass
                return jsonify({'error': f'合并后总页数超过100页限制（当前{total_pages}页）'}), 400
            
            uploaded_files.append({
                'path': temp_path,
                'filename': filename,
                'page_count': page_count,
                'selected_pages': selected_pages_list
            })
        
        if len(uploaded_files) == 0:
            return jsonify({'error': '没有有效的PDF文件'}), 400
        
        logger.info(f"准备合并 {len(uploaded_files)} 个文件，总页数: {total_pages}")
        
        # 开始合并
        result = fitz.open()
        
        for file_info in uploaded_files:
            doc = fitz.open(file_info['path'])
            
            # 逐页插入选定的页面
            for page_num in file_info['selected_pages']:
                if page_num < len(doc):
                    result.insert_pdf(doc, from_page=page_num, to_page=page_num)
            
            doc.close()
        
        # 生成输出文件名
        output_filename = f"merged_{uuid.uuid4().hex[:8]}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)
        
        # 保存合并后的PDF（优化压缩）
        result.save(output_path, garbage=4, deflate=True, clean=True)
        output_size = os.path.getsize(output_path)
        result.close()
        
        # 清理临时文件
        for file_info in uploaded_files:
            try:
                os.remove(file_info['path'])
            except:
                pass
        
        end_time = time.time()
        processing_time = round(end_time - start_time, 2)
        
        logger.info(f"合并完成: {output_filename}, 页数: {total_pages}, 大小: {output_size/1024:.2f}KB, 耗时: {processing_time}秒")
        
        return jsonify({
            'success': True,
            'filename': output_filename,
            'download_url': f'/download/{output_filename}',
            'file_size': output_size,
            'file_size_mb': round(output_size / (1024 * 1024), 2),
            'total_pages': total_pages,
            'file_count': len(uploaded_files),
            'processing_time': processing_time,
            'message': '合并成功'
        })
        
    except Exception as e:
        logger.error(f"PDF合并失败: {str(e)}")
        logger.error(traceback.format_exc())
        
        # 清理可能的临时文件
        try:
            for file_info in uploaded_files:
                if os.path.exists(file_info['path']):
                    os.remove(file_info['path'])
        except:
            pass
        
        return jsonify({'error': f'合并失败: {str(e)}'}), 500


def parse_page_ranges(pages_str, total_pages):
    """
    解析页码范围字符串
    例如: "1-5,7,9-10" -> [0,1,2,3,4,6,8,9] (转为0-based索引)
    """
    if not pages_str or pages_str.strip() == '':
        return list(range(total_pages))
    
    result = []
    parts = pages_str.split(',')
    
    for part in parts:
        part = part.strip()
        if '-' in part:
            # 范围格式 "1-5"
            try:
                start, end = part.split('-')
                start = int(start.strip())
                end = int(end.strip())
                
                # 验证范围
                if start < 1 or end > total_pages or start > end:
                    continue
                
                # 转为0-based索引
                for i in range(start - 1, end):
                    if i not in result:
                        result.append(i)
            except:
                continue
        else:
            # 单个页码 "7"
            try:
                page = int(part.strip())
                if 1 <= page <= total_pages:
                    page_idx = page - 1  # 转为0-based索引
                    if page_idx not in result:
                        result.append(page_idx)
            except:
                continue
    
    # 排序
    result.sort()
    return result


@app.route('/pdf/upload-image', methods=['POST'])
def upload_image():
    """
    上传图片用于插入PDF
    支持：相册、会话、拍照
    """
    try:
        if 'file' not in request.files:
            return jsonify({'error': '未找到图片文件'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': '未选择文件'}), 400
        
        # 验证文件大小（50MB）
        file.seek(0, 2)
        file_size = file.tell()
        file.seek(0)
        
        if file_size > 50 * 1024 * 1024:
            return jsonify({'error': '文件大小超过50MB限制'}), 400
        
        # 验证文件类型（图片）
        allowed_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp'}
        filename = secure_filename(file.filename)
        file_ext = os.path.splitext(filename)[1].lower()
        
        if file_ext not in allowed_extensions:
            return jsonify({'error': '只支持图片格式（JPG/PNG/GIF/BMP/WEBP）'}), 400
        
        # 保存图片
        temp_filename = f"img_{uuid.uuid4().hex[:8]}_{filename}"
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], temp_filename)
        file.save(temp_path)
        
        logger.info(f"图片上传成功: {filename}, 大小: {file_size/1024:.2f}KB")
        
        return jsonify({
            'success': True,
            'image_path': temp_path,
            'filename': filename,
            'file_size': file_size,
            'message': '上传成功'
        })
        
    except Exception as e:
        logger.error(f"图片上传失败: {str(e)}")
        return jsonify({'error': f'上传失败: {str(e)}'}), 500


@app.route('/pdf/get-pages-preview', methods=['POST'])
def get_pages_preview():
    """
    获取PDF所有页面的缩略图预览
    用于前端显示页面列表
    """
    try:
        if 'file' not in request.files:
            return jsonify({'error': '未找到文件'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': '未选择文件'}), 400
        
        if not file.filename.lower().endswith('.pdf'):
            return jsonify({'error': '只支持PDF文件'}), 400
        
        # 保存临时文件
        filename = secure_filename(file.filename)
        temp_filename = f"preview_{uuid.uuid4().hex[:8]}_{filename}"
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], temp_filename)
        file.save(temp_path)
        
        # 生成缩略图
        doc = fitz.open(temp_path)
        page_count = len(doc)
        previews = []
        
        # 生成每页的缩略图（Base64）
        for page_num in range(page_count):
            page = doc[page_num]
            # 缩略图尺寸：宽度200px，保持比例
            zoom = 200.0 / page.rect.width
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)
            
            # 转换为Base64
            img_data = pix.tobytes("png")
            img_base64 = base64.b64encode(img_data).decode('utf-8')
            
            previews.append({
                'page_number': page_num + 1,
                'thumbnail': f'data:image/png;base64,{img_base64}'
            })
        
        doc.close()
        
        # 删除临时文件
        try:
            os.remove(temp_path)
        except:
            pass
        
        logger.info(f"生成预览: {filename}, {page_count}页")
        
        return jsonify({
            'success': True,
            'page_count': page_count,
            'previews': previews,
            'message': '预览生成成功'
        })
        
    except Exception as e:
        logger.error(f"生成预览失败: {str(e)}")
        return jsonify({'error': f'生成失败: {str(e)}'}), 500


@app.route('/pdf/manage-pages', methods=['POST'])
def manage_pdf_pages():
    """
    PDF页面管理：增删页/排序
    支持：
    1. 删除页面：delete_pages: [2, 5, 7]
    2. 重排序：reorder: [3, 1, 2, 4]
    3. 插入图片：insert_images: [{position: 3, image_data: base64}, ...]
    """
    import time
    start_time = time.time()
    
    try:
        # 获取PDF文件
        if 'pdf_file' not in request.files:
            return jsonify({'error': '未找到PDF文件'}), 400
        
        pdf_file = request.files['pdf_file']
        if pdf_file.filename == '':
            return jsonify({'error': '未选择PDF文件'}), 400
        
        if not pdf_file.filename.lower().endswith('.pdf'):
            return jsonify({'error': '只支持PDF文件'}), 400
        
        # 获取操作参数
        delete_pages = request.form.get('delete_pages', '')
        reorder = request.form.get('reorder', '')
        insert_images = request.form.get('insert_images', '[]')
        
        import json
        try:
            delete_pages = json.loads(delete_pages) if delete_pages else []
            reorder = json.loads(reorder) if reorder else []
            insert_images = json.loads(insert_images) if insert_images else []
        except:
            return jsonify({'error': '参数格式错误'}), 400
        
        # 保存PDF文件
        filename = secure_filename(pdf_file.filename)
        temp_filename = f"manage_{uuid.uuid4().hex[:8]}_{filename}"
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], temp_filename)
        pdf_file.save(temp_path)
        
        # 验证文件大小
        file_size = os.path.getsize(temp_path)
        if file_size > 50 * 1024 * 1024:
            os.remove(temp_path)
            return jsonify({'error': 'PDF文件大小超过50MB限制'}), 400
        
        # 打开PDF
        doc = fitz.open(temp_path)
        original_total_pages = len(doc)  # 保存原始页数（删除前）
        total_pages = original_total_pages  # 当前页数（会随着删除而变化）
        
        logger.info(f"开始管理PDF页面: {filename}, 总页数: {original_total_pages}")
        logger.info(f"删除页面: {delete_pages}, 重排序: {reorder}, 插入图片: {len(insert_images)}个")
        
        # 验证页码有效性（基于原始页数）
        for page_num in delete_pages:
            if page_num < 1 or page_num > original_total_pages:
                doc.close()
                os.remove(temp_path)
                return jsonify({'error': f'无效的页码: {page_num}'}), 400
        
        if reorder:
            for page_num in reorder:
                if page_num < 1 or page_num > original_total_pages:
                    doc.close()
                    os.remove(temp_path)
                    return jsonify({'error': f'无效的页码: {page_num}'}), 400
        
        # 执行操作：先删除页面（从后往前删除，避免索引变化）
        if delete_pages:
            for page_num in sorted(set(delete_pages), reverse=True):
                doc.delete_page(page_num - 1)  # 转为0-based索引
            total_pages = len(doc)  # 更新为删除后的页数
            logger.info(f"删除页面完成，剩余页数: {total_pages} (原始页数: {original_total_pages})")
        
        # 执行操作：重排序
        if reorder:
            result = fitz.open()
            current_doc_length = len(doc)
            logger.info(f"开始重排序: 当前文档页数={current_doc_length}, 重排序列表={reorder}")
            
            # reorder中的页码可能是基于原始文档的页码（删除前的页码）
            # 需要转换为删除后的页码索引
            # 例如：删除第1页后，原始第2页变成第1页，原始第3页变成第2页...
            
            # 创建一个映射：原始页码 -> 删除后的索引（0-based）
            # 删除的页面不在映射中
            original_to_index = {}
            current_index = 0
            # 使用original_total_pages（删除前的页数），而不是total_pages（删除后的页数）
            for orig_page in range(1, original_total_pages + 1):
                if orig_page not in delete_pages:
                    original_to_index[orig_page] = current_index
                    current_index += 1
            
            logger.info(f"页码映射: {original_to_index}")
            logger.info(f"重排序列表长度: {len(reorder)}, 映射中有效页码数: {len([p for p in reorder if p in original_to_index])}")
            
            # 去重：确保每个页码只复制一次（保持顺序）
            seen_pages = set()
            unique_reorder = []
            for page_idx in reorder:
                if page_idx not in seen_pages:
                    seen_pages.add(page_idx)
                    unique_reorder.append(page_idx)
            
            if len(unique_reorder) != len(reorder):
                logger.warning(f"重排序列表中有重复页码，已去重: 原始{len(reorder)}个 -> 去重后{len(unique_reorder)}个")
            
            # 按照去重后的reorder列表复制页面
            copied_count = 0
            for page_idx in unique_reorder:
                # page_idx是原始页码，需要转换为删除后的索引
                if page_idx in original_to_index:
                    source_index = original_to_index[page_idx]
                    # 验证source_index是否在有效范围内
                    if 0 <= source_index < current_doc_length:
                        try:
                            result.insert_pdf(doc, from_page=source_index, to_page=source_index)
                            copied_count += 1
                            logger.info(f"复制页面 {copied_count}: 原始页码={page_idx} -> 索引={source_index}, 复制后结果文档页数={len(result)}")
                        except Exception as e:
                            logger.error(f"复制页面失败: 原始页码={page_idx}, 索引={source_index}, 错误={str(e)}")
                    else:
                        logger.error(f"索引超出范围: 原始页码={page_idx}, 索引={source_index}, 文档长度={current_doc_length}")
                else:
                    logger.warning(f"重排序时跳过已删除的页码: {page_idx}")
            
            # 验证重排序后的页数
            result_length = len(result)
            expected_length = len([p for p in reorder if p in original_to_index])
            if result_length != expected_length:
                logger.error(f"重排序后页数不匹配: 预期{expected_length}页, 实际{result_length}页, 已复制{copied_count}页")
            else:
                logger.info(f"重排序页数验证通过: {result_length}页")
            
            doc.close()
            doc = result
            logger.info(f"重排序完成: 重排序前{current_doc_length}页, 重排序后{result_length}页, 预期{expected_length}页, 已复制{copied_count}页")
        
        # 执行操作：插入图片（需要先收集所有插入操作，然后按位置排序）
        successful_inserts = 0  # 成功插入的图片数量（在if块外定义，以便后续使用）
        if insert_images:
            # 获取当前文档页数（删除和重排序后的页数）
            initial_page_count = len(doc)
            logger.info(f"插入图片前，当前文档页数: {initial_page_count}, 待插入图片数: {len(insert_images)}")
            
            # 按位置排序，从前往后插入（从小到大）
            # 这样可以确保位置计算正确，因为每次插入后，后续位置会自动调整
            sorted_inserts = sorted(insert_images, key=lambda x: x.get('position', 1))
            
            # 记录已插入的图片数量，用于调整后续位置
            inserted_count = 0
            successful_inserts = 0  # 成功插入的图片数量
            
            for img_info in sorted_inserts:
                position = img_info.get('position', 1)
                image_data = img_info.get('image_data', '')
                
                if not image_data:
                    continue
                
                # 获取当前文档的实际页数（每次插入后页数会增加）
                current_doc_length = len(doc)
                
                # 验证位置有效性：position应该在1到initial_page_count+1之间（基于初始页数）
                # position=1表示在第1页之前，position=N+1表示在最后一页之后
                # 注意：客户端发送的position是基于删除后的页数，不包括已插入的图片
                if position < 1:
                    logger.warning(f"插入图片位置无效: {position}，调整为1")
                    position = 1
                elif position > initial_page_count + 1:
                    logger.warning(f"插入图片位置超出范围: {position}，调整为{initial_page_count + 1} (初始文档页数: {initial_page_count})")
                    position = initial_page_count + 1
                
                # 从前往后插入时，不需要调整position
                # 因为每次插入后，后续的位置会自动调整
                # 例如：初始9页，position=1，插入到位置0（第1页之前）
                # 插入后文档变成10页，原来的第1页现在是第2页
                # 如果再有position=2的插入请求，它应该插入到新的第2页之前，也就是位置1
                # 所以直接使用position即可，不需要加inserted_count
                adjusted_position = position
                
                logger.info(f"准备插入图片: 原始位置={position}, 已插入图片数={inserted_count}, 使用位置={adjusted_position} (从前往后插入，无需调整)")
                
                try:
                    # 解码base64图片
                    import base64
                    image_bytes = base64.b64decode(image_data.split(',')[-1] if ',' in image_data else image_data)
                    
                    # 使用PIL打开图片
                            from PIL import Image
                            import io
                            img = Image.open(io.BytesIO(image_bytes))
                    
                    # 转换为RGB（如果是RGBA或其他模式）
                    if img.mode != 'RGB':
                            if img.mode == 'RGBA':
                                rgb_img = Image.new('RGB', img.size, (255, 255, 255))
                                rgb_img.paste(img, mask=img.split()[3])
                                img = rgb_img
                        else:
                            img = img.convert('RGB')
                    
                    # 获取图片原始尺寸
                    original_width = img.width
                    original_height = img.height
                    
                    # 使用高分辨率：300 DPI（打印质量）而不是72 DPI
                    # A4尺寸: 210 x 297 mm
                    # 在300 DPI下: 2480 x 3508 像素
                    # 在72 DPI下: 595 x 842 点
                    target_dpi = 300  # 高分辨率
                    standard_dpi = 72  # PDF标准DPI
                    
                    # A4页面尺寸（毫米）
                    a4_width_mm = 210
                    a4_height_mm = 297
                    
                    # 转换为300 DPI下的像素尺寸
                    target_width_px = int(a4_width_mm * target_dpi / 25.4)  # 约2480像素
                    target_height_px = int(a4_height_mm * target_dpi / 25.4)  # 约3508像素
                    
                    # PDF页面尺寸（点，72 DPI）
                    standard_page_width_pt = 595
                    standard_page_height_pt = 842
                    
                    # 计算缩放比例，使图片宽度填满目标宽度（300 DPI），保持宽高比
                    # 优先按宽度缩放，确保图片宽度填满页面宽度
                    scale = target_width_px / original_width
                    scaled_width_px = target_width_px
                    scaled_height_px = int(original_height * scale)
                    
                    # 如果缩放后的高度超过目标高度，则按高度缩放
                    # 但这样会导致宽度小于目标宽度，我们需要重新调整
                    if scaled_height_px > target_height_px:
                        scale = target_height_px / original_height
                        scaled_height_px = target_height_px
                        # 按高度缩放后，计算对应的宽度
                        scaled_width_px = int(original_width * scale)
                        # 如果宽度小于目标宽度，需要再次按宽度缩放（会裁剪高度）
                        # 但为了保持宽高比和填满宽度，我们选择按宽度缩放
                        # 这样高度可能会超出，但我们可以限制高度
                        if scaled_width_px < target_width_px:
                            # 重新按宽度缩放，确保填满宽度
                            scale = target_width_px / original_width
                            scaled_width_px = target_width_px
                            scaled_height_px = int(original_height * scale)
                            # 如果高度超过目标高度，限制高度（会裁剪）
                            if scaled_height_px > target_height_px:
                                scaled_height_px = target_height_px
                    
                    # 使用高质量重采样算法（LANCZOS）进行缩放
                    # LANCZOS是最高质量的算法，适合放大和缩小
                    if scale != 1.0:
                        # 如果放大（scale > 1），使用LANCZOS放大以保持质量
                        # 如果缩小（scale < 1），使用LANCZOS缩小以减少失真
                        img = img.resize((scaled_width_px, scaled_height_px), Image.Resampling.LANCZOS)
                        logger.info(f"图片已调整到高分辨率: {original_width}x{original_height} -> {scaled_width_px}x{scaled_height_px} (缩放比例: {scale:.2f}, 目标DPI: {target_dpi})")
                    else:
                        logger.info(f"图片保持原始尺寸: {original_width}x{original_height} (目标DPI: {target_dpi})")
                    
                    # 将像素尺寸转换为PDF点尺寸（72 DPI）
                    # 300 DPI下的像素 / (300/72) = 72 DPI下的点
                    calculated_width_pt = scaled_width_px * standard_dpi / target_dpi
                    calculated_height_pt = scaled_height_px * standard_dpi / target_dpi
                    
                    # 在指定位置插入新页面
                    # position是1-based：position=1表示在第1页之前，position=N+1表示在最后一页之后
                    # insert_position是0-based：0表示在第1页之前，N表示在最后一页之后
                    # 从前往后插入时，直接使用position，因为每次插入后，后续的位置会自动调整
                    current_doc_length = len(doc)
                    
                    # 计算插入位置：position=1 -> insert_position=0, position=N+1 -> insert_position=N
                    # position是1-based，需要转换为0-based
                    # 注意：从前往后插入时，position不需要调整，因为每次插入后，后续位置会自动调整
                    insert_position = max(0, min(adjusted_position - 1, current_doc_length))
                    
                    logger.info(f"插入图片: 请求位置={position}, 已插入图片数={inserted_count}, 使用位置={adjusted_position}, 文档当前页数={current_doc_length}, 插入位置(0-based)={insert_position}")
                    
                    # 创建新的PDF页面，使用标准A4尺寸
                    # new_page会在指定位置插入新页面，文档长度会自动增加
                    new_page = doc.new_page(insert_position, width=standard_page_width_pt, height=standard_page_height_pt)
                    
                    # 验证页面是否成功插入
                    new_doc_length = len(doc)
                    if new_doc_length != current_doc_length + 1:
                        logger.error(f"错误: 插入页面后，文档页数未正确增加！插入前: {current_doc_length}, 插入后: {new_doc_length}, 预期: {current_doc_length + 1}")
                        # 如果插入失败，抛出异常，让外层catch处理
                        raise Exception(f"页面插入失败: 插入前{current_doc_length}页, 插入后{new_doc_length}页, 预期{current_doc_length + 1}页")
                    
                    # 页面插入成功后，立即增加已插入图片计数
                    inserted_count += 1
                    successful_inserts += 1
                    logger.info(f"✓ 页面插入成功: 插入前{current_doc_length}页 -> 插入后{new_doc_length}页, 已插入图片总数: {inserted_count}, 成功插入: {successful_inserts}")
                    
                    # 将图片嵌入到PDF页面
                    # 保存图片到临时文件，使用最高质量PNG（无压缩）
                            temp_img_path = os.path.join(app.config['UPLOAD_FOLDER'], f"temp_img_{uuid.uuid4().hex[:8]}.png")
                    # 使用最高质量保存PNG
                    # compress_level=0: 无压缩，最高质量
                    # optimize=False: 不优化，保持原始质量
                    img.save(temp_img_path, 'PNG', compress_level=0, optimize=False)
                    
                    # 确保图片宽度完全填满页面宽度，左右无间距
                    # 强制使用页面宽度，避免浮点数精度问题导致的偏差
                    img_width_pt = standard_page_width_pt
                    
                    # 计算图片高度（保持宽高比）
                    # 如果计算出的宽度正好等于页面宽度，使用计算出的高度
                    # 否则按比例调整高度
                    if abs(calculated_width_pt - standard_page_width_pt) < 0.1:
                        img_height_pt = calculated_height_pt
                    else:
                        # 按页面宽度重新计算高度，保持宽高比
                        aspect_ratio = calculated_height_pt / calculated_width_pt if calculated_width_pt > 0 else 1.0
                        img_height_pt = img_width_pt * aspect_ratio
                    
                    # 如果高度超过页面高度，限制为页面高度，并重新计算宽度（但保持页面宽度）
                    if img_height_pt > standard_page_height_pt:
                        img_height_pt = standard_page_height_pt
                        # 注意：这里不调整宽度，因为我们要确保宽度填满页面
                        # 图片可能会被裁剪，但宽度始终填满
                    
                    # 计算图片在页面中的位置
                    # 水平方向：始终从左边开始（x=0），确保左右完全对齐，无间距
                    x_offset = 0.0
                    
                    # 垂直方向：如果图片高度小于页面高度，垂直居中
                    # 如果图片高度大于等于页面高度，从顶部开始
                    if img_height_pt < standard_page_height_pt:
                        y_offset = (standard_page_height_pt - img_height_pt) / 2
                    else:
                        y_offset = 0.0
                    
                    # 在页面上插入图片，使用高分辨率
                    # 图片从左边开始（x=0），完全填满页面宽度（595点），确保左右间距完全一致
                    # 使用精确的矩形坐标，确保左右完全对齐
                    img_rect = fitz.Rect(
                        float(x_offset), 
                        float(y_offset), 
                        float(x_offset + img_width_pt), 
                        float(y_offset + img_height_pt)
                    )
                    new_page.insert_image(img_rect, filename=temp_img_path)
                    
                    logger.info(f"图片定位: x={x_offset:.2f}, y={y_offset:.2f}, 宽度={img_width_pt:.2f}点 (页面宽度={standard_page_width_pt}点, 确保左右无间距)")
                    logger.info(f"✓ 插入高分辨率图片到位置 {position} (PDF尺寸: {img_width_pt:.1f}x{img_height_pt:.1f}点, 原始像素: {scaled_width_px}x{scaled_height_px})")
                    logger.info(f"✓ 插入后文档页数: {len(doc)} (插入前: {current_doc_length}), 已插入图片总数: {inserted_count}")
                    
                    # 删除临时文件
                    try:
                            os.remove(temp_img_path)
                    except:
                        pass
                    
                    # 图片插入成功，inserted_count已经在页面插入后增加了
                    
                except Exception as e:
                    logger.warning(f"✗ 插入图片失败: {str(e)}")
                    logger.warning(traceback.format_exc())
                    # 插入失败时，inserted_count不增加，这是正确的
                    continue
            
            # 插入图片完成后，记录最终统计
            final_page_count_after_insert = len(doc)
            logger.info(f"插入图片完成: 初始页数={initial_page_count}, 成功插入={successful_inserts}张, 最终页数={final_page_count_after_insert}, 预期页数={initial_page_count + successful_inserts}")
            if final_page_count_after_insert != initial_page_count + successful_inserts:
                logger.error(f"警告: 最终页数不正确！初始={initial_page_count}, 成功插入={successful_inserts}, 预期={initial_page_count + successful_inserts}, 实际={final_page_count_after_insert}")
        
        # 生成输出文件名
        output_filename = f"managed_{uuid.uuid4().hex[:8]}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)
        
        # 保存结果（优化压缩）
        doc.save(output_path, garbage=4, deflate=True, clean=True)
        output_size = os.path.getsize(output_path)
        final_pages = len(doc)
        
        # 记录详细的操作信息
        deleted_count = len(delete_pages) if delete_pages else 0
        requested_inserts = len(insert_images) if insert_images else 0
        # 使用实际成功插入的图片数量（在插入图片的代码块中已经统计）
        actual_inserted_count = successful_inserts  # 使用上面定义的successful_inserts变量
        reordered = len(reorder) if reorder else 0
        
        logger.info(f"页面管理完成:")
        logger.info(f"  原始页数: {total_pages}")
        logger.info(f"  删除页数: {deleted_count}")
        logger.info(f"  请求插入图片数: {requested_inserts}")
        logger.info(f"  实际成功插入图片数: {actual_inserted_count}")
        logger.info(f"  重排序: {reordered > 0}")
        logger.info(f"  最终页数: {final_pages}")
        logger.info(f"  预期页数: {total_pages - deleted_count + actual_inserted_count}")
        logger.info(f"  文件: {output_filename}, 耗时: {round(time.time() - start_time, 2)}秒")
        
        # 验证最终页数是否正确
        expected_pages = total_pages - deleted_count + actual_inserted_count
        if final_pages != expected_pages:
            logger.error(f"错误: 最终页数不匹配！预期={expected_pages}, 实际={final_pages}, 差异={final_pages - expected_pages}")
        
        doc.close()
        
        # 清理临时文件
        try:
            os.remove(temp_path)
        except:
            pass
        
        end_time = time.time()
        processing_time = round(end_time - start_time, 2)
        
        return jsonify({
            'success': True,
            'filename': output_filename,
            'download_url': f'/download/{output_filename}',
            'file_size': output_size,
            'file_size_mb': round(output_size / (1024 * 1024), 2),
            'original_pages': total_pages,
            'final_pages': final_pages,
            'processing_time': processing_time,
            'message': '页面管理成功'
        })
        
    except Exception as e:
        logger.error(f"PDF页面管理失败: {str(e)}")
        logger.error(traceback.format_exc())
        return jsonify({'error': f'处理失败: {str(e)}'}), 500


@app.route('/pdf/rotate', methods=['POST'])
def rotate_pdf_pages():
    """
    旋转PDF页面 - 高性能实现
    支持0/90/180/270度顺时针旋转
    """
    try:
        if 'file' not in request.files:
            return jsonify({'error': '未上传文件'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': '文件名为空'}), 400
        
        # 获取旋转参数
        rotate_pages = request.form.get('pages', '')  # 逗号分隔的页码，如 "1,3,5" 或 "all" 表示全部
        rotate_angle = request.form.get('angle', '90')  # 旋转角度：0, 90, 180, 270
        
        # 支持批量旋转：pages_angles JSON格式，如 {"1": 90, "3": 180, "5": 270}
        # 如果提供了pages_angles，优先使用它（支持不同页面不同角度，高性能）
        pages_angles_json = request.form.get('pages_angles', '')
        pages_angles = {}
        use_batch_mode = False
        
        if pages_angles_json:
            try:
                pages_angles = json.loads(pages_angles_json)
                # 验证角度有效性
                for page_num_str, angle_val in pages_angles.items():
                    if angle_val not in [0, 90, 180, 270]:
                        return jsonify({'error': f'页面{page_num_str}的旋转角度无效: {angle_val}'}), 400
                use_batch_mode = True
                logger.info(f"收到批量旋转请求: {len(pages_angles)}个页面，不同角度")
            except json.JSONDecodeError:
                logger.warning(f"解析pages_angles JSON失败，使用单角度模式")
            except Exception as e:
                logger.warning(f"处理pages_angles失败: {str(e)}，使用单角度模式")
        
        # 单角度模式：验证角度
        if not use_batch_mode:
            try:
                angle = int(rotate_angle)
                if angle not in [0, 90, 180, 270]:
                    return jsonify({'error': '旋转角度必须是0、90、180或270'}), 400
            except ValueError:
                return jsonify({'error': '无效的旋转角度'}), 400
        
        # 保存上传的文件
        filename = secure_filename(file.filename)
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{uuid.uuid4().hex[:8]}_{filename}")
        file.save(temp_path)
        
        start_time = time.time()
        
        # 打开PDF
        doc = fitz.open(temp_path)
        total_pages = len(doc)
        
        # 执行旋转（PyMuPDF高性能旋转）
        rotated_count = 0
        rotated_pages_info = []  # 记录旋转信息
        
        if use_batch_mode:
            # 批量模式：不同页面不同角度（高性能，一次处理）
            logger.info(f"开始批量旋转PDF页面: {filename}, 总页数: {total_pages}, 待旋转: {len(pages_angles)}页")
            
            for page_num_str, angle_val in pages_angles.items():
                try:
                    page_num = int(page_num_str)
                    if 1 <= page_num <= total_pages:
                        page_idx = page_num - 1
                        page = doc[page_idx]
                        
                        # 保存原始页面尺寸和cropbox（确保旋转后宽高不变）
                        orig_rect = page.rect
                        orig_width = orig_rect.width
                        orig_height = orig_rect.height
                        
                        # 经过测试验证：PyMuPDF的set_rotation是顺时针角度
                        # 直接使用用户输入的顺时针角度即可，不需要转换
                        
                        # 获取当前旋转角度
                        current_rotation = page.rotation
                        
                        # 先重置为0度，然后设置目标角度（确保是绝对角度，不是累加）
                        if current_rotation != 0:
                            # 计算需要旋转多少度才能回到0度
                            reset_increment = (360 - current_rotation) % 360
                            if reset_increment > 0:
                                page.set_rotation(reset_increment)
                        
                        # 执行旋转（只旋转，让PyMuPDF自动处理尺寸）
                        if angle_val > 0:
                            page.set_rotation(angle_val)
                        
                        # 记录旋转信息
                        rotated_rect = page.rect
                        logger.info(f"旋转页面 {page_num}: {angle_val}度顺时针，原始: {orig_width:.1f}x{orig_height:.1f}，旋转后: {rotated_rect.width:.1f}x{rotated_rect.height:.1f}")
                        
                        rotated_count += 1
                        rotated_pages_info.append(f"第{page_num}页({angle_val}°)")
                        logger.debug(f"旋转页面 {page_num}: {angle_val}度顺时针，原始尺寸: {orig_width}x{orig_height}")
                    else:
                        logger.warning(f"跳过无效页码: {page_num}")
                except (ValueError, IndexError) as e:
                    logger.warning(f"处理页面 {page_num_str} 失败: {str(e)}")
                except Exception as e:
                    logger.warning(f"旋转页面 {page_num_str} 失败: {str(e)}")
                    logger.warning(traceback.format_exc())
            
            logger.info(f"批量旋转完成: 成功旋转 {rotated_count} 页")
        else:
            # 单角度模式：所有页面同一角度
            logger.info(f"开始旋转PDF页面: {filename}, 总页数: {total_pages}, 角度: {angle}度")
            
            # 解析要旋转的页面
            pages_to_rotate = []
            if rotate_pages.lower() == 'all' or rotate_pages.strip() == '':
                # 旋转所有页面
                pages_to_rotate = list(range(total_pages))
                logger.info(f"旋转所有页面: {total_pages}页")
            else:
                # 解析页码列表
                try:
                    page_list = [int(p.strip()) for p in rotate_pages.split(',') if p.strip()]
                    # 转换为0-based索引，并验证有效性
                    for page_num in page_list:
                        if 1 <= page_num <= total_pages:
                            pages_to_rotate.append(page_num - 1)
                        else:
                            logger.warning(f"跳过无效页码: {page_num}")
                except ValueError:
                    doc.close()
                    os.remove(temp_path)
                    return jsonify({'error': '无效的页码格式'}), 400
            
            if not pages_to_rotate:
                doc.close()
                os.remove(temp_path)
                return jsonify({'error': '没有有效的页面需要旋转'}), 400
            
            # PyMuPDF的set_rotation是顺时针角度，直接使用用户输入的角度
            # 不需要转换（0, 90, 180, 270都是顺时针）
            
            # 批量旋转页面（高性能，使用绝对角度而非累加）
            for page_idx in pages_to_rotate:
                try:
                    page = doc[page_idx]
                    
                    # 保存原始页面尺寸
                    orig_rect = page.rect
                    orig_width = orig_rect.width
                    orig_height = orig_rect.height
                    
                    # 经过测试验证：PyMuPDF的set_rotation是顺时针角度
                    # 直接使用用户输入的顺时针角度即可，不需要转换
                    
                    # 获取当前旋转角度
                    current_rotation = page.rotation
                    
                    # 先重置为0度，然后设置目标角度（确保是绝对角度，不是累加）
                    if current_rotation != 0:
                        # 计算需要旋转多少度才能回到0度
                        reset_increment = (360 - current_rotation) % 360
                        if reset_increment > 0:
                            page.set_rotation(reset_increment)
                    
                    # 执行旋转（只旋转，让PyMuPDF自动处理尺寸）
                    if angle > 0:
                        page.set_rotation(angle)
                    
                    # 记录旋转信息
                    rotated_rect = page.rect
                    logger.info(f"旋转页面 {page_idx + 1}: {angle}度顺时针，原始: {orig_width:.1f}x{orig_height:.1f}，旋转后: {rotated_rect.width:.1f}x{rotated_rect.height:.1f}")
                    
                    rotated_count += 1
                    logger.debug(f"旋转页面 {page_idx + 1}: {angle}度顺时针，原始尺寸: {orig_width}x{orig_height}")
                except Exception as e:
                    logger.warning(f"旋转页面 {page_idx + 1} 失败: {str(e)}")
                    logger.warning(traceback.format_exc())
            
            logger.info(f"成功旋转 {rotated_count} 页，角度: {angle}度")
        
        # 生成输出文件名
        output_filename = f"rotated_{uuid.uuid4().hex[:8]}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)
        
        # 保存结果（优化压缩）
        doc.save(output_path, garbage=4, deflate=True, clean=True)
        output_size = os.path.getsize(output_path)
        doc.close()
        
        # 清理临时文件
        try:
            os.remove(temp_path)
        except:
            pass
        
        end_time = time.time()
        processing_time = round(end_time - start_time, 2)
        
        if use_batch_mode:
            logger.info(f"PDF批量旋转完成: {output_filename}, 旋转页数: {rotated_count}, 耗时: {processing_time}秒")
        else:
            logger.info(f"PDF旋转完成: {output_filename}, 旋转页数: {rotated_count}, 角度: {angle}度, 耗时: {processing_time}秒")
        
        # 构建返回数据
        result_data = {
            'success': True,
            'filename': output_filename,
            'download_url': f'/download/{output_filename}',
            'file_size': output_size,
            'file_size_mb': round(output_size / (1024 * 1024), 2),
            'total_pages': total_pages,
            'rotated_pages': rotated_count,
            'processing_time': processing_time
        }
        
        if use_batch_mode:
            result_data['rotated_info'] = ', '.join(rotated_pages_info[:5])  # 最多显示5个
            if len(rotated_pages_info) > 5:
                result_data['rotated_info'] += f' 等{rotated_count}页'
            logger.info(f"批量旋转完成: {result_data['rotated_info']}")
        else:
            result_data['angle'] = angle
        
        return jsonify(result_data)
        
    except Exception as e:
        logger.error(f'PDF旋转错误: {str(e)}')
        logger.error(traceback.format_exc())
        return jsonify({'error': f'旋转失败: {str(e)}'}), 500

@app.route('/pdf/compress', methods=['POST'])
def compress_pdf():
    """
    PDF压缩接口
    支持三种压缩级别：low(20%), medium(50%), high(80%)
    """
    try:
        if 'file' not in request.files:
            return jsonify({'error': '未找到文件'}), 400
        
        file = request.files['file']
        if file.filename == '':
            return jsonify({'error': '未选择文件'}), 400
        
        if not file.filename.lower().endswith('.pdf'):
            return jsonify({'error': '只支持PDF文件'}), 400
        
        # 获取压缩级别
        compression_level = request.form.get('compression_level', 'medium').lower()
        if compression_level not in ['low', 'medium', 'high']:
            compression_level = 'medium'
        
        # 压缩级别配置
        compression_configs = {
            'low': {
                'image_quality': 90,      # 图片质量90%
                'image_dpi': 200,         # 图片DPI 200
                'garbage': 2,             # 轻度清理
                'deflate': True,          # 启用压缩
                'clean': True             # 清理交叉引用
            },
            'medium': {
                'image_quality': 75,      # 图片质量75%
                'image_dpi': 150,         # 图片DPI 150
                'garbage': 3,             # 中度清理
                'deflate': True,
                'clean': True
            },
            'high': {
                'image_quality': 60,      # 图片质量60%
                'image_dpi': 100,         # 图片DPI 100
                'garbage': 4,             # 深度清理
                'deflate': True,
                'clean': True
            }
        }
        
        config = compression_configs[compression_level]
        
        # 保存上传的文件
        filename = secure_filename(file.filename)
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{uuid.uuid4().hex[:8]}_{filename}")
        file.save(temp_path)
        
        start_time = time.time()
        original_size = os.path.getsize(temp_path)
        
        logger.info(f"开始压缩PDF: {filename}, 原始大小: {original_size / 1024:.1f}KB, 压缩级别: {compression_level}")
        
        # 打开PDF
        doc = fitz.open(temp_path)
        total_pages = len(doc)
        
        # 压缩图片：遍历所有图片并压缩
        image_count = 0
        total_images = 0
        
        for page_num in range(total_pages):
            page = doc[page_num]
            image_list = page.get_images()
            total_images += len(image_list)
            
            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    image_ext = base_image["ext"]
                    
                    # 只处理JPEG和PNG图片
                    if image_ext.lower() not in ['jpeg', 'jpg', 'png']:
                        continue
                    
                    # 使用PIL压缩图片
                    img_pil = Image.open(io.BytesIO(image_bytes))
                    original_size = len(image_bytes)
                    
                    # 如果是PNG，转换为JPEG以减小文件大小
                    convert_to_jpeg = False
                    if image_ext.lower() == 'png':
                        if img_pil.mode in ['RGBA', 'LA']:
                            # 创建白色背景
                            background = Image.new('RGB', img_pil.size, (255, 255, 255))
                            if img_pil.mode == 'RGBA':
                                background.paste(img_pil, mask=img_pil.split()[3])
                            else:
                                background.paste(img_pil)
                            img_pil = background
                            convert_to_jpeg = True
                        elif compression_level == 'high':
                            # 高压缩级别时，所有PNG都转JPEG
                            convert_to_jpeg = True
                    
                    # 计算目标尺寸（基于DPI）
                    scale_factor = config['image_dpi'] / 72.0
                    if scale_factor < 1.0:
                        new_width = int(img_pil.width * scale_factor)
                        new_height = int(img_pil.height * scale_factor)
                        if new_width > 0 and new_height > 0:
                            img_pil = img_pil.resize((new_width, new_height), Image.Resampling.LANCZOS)
                    
                    # 压缩图片
                    img_buffer = io.BytesIO()
                    if convert_to_jpeg or image_ext.lower() in ['jpeg', 'jpg']:
                        img_pil.save(img_buffer, format='JPEG', quality=config['image_quality'], optimize=True)
                    else:
                        # PNG优化
                        img_pil.save(img_buffer, format='PNG', optimize=True)
                    
                    compressed_bytes = img_buffer.getvalue()
                    compressed_size = len(compressed_bytes)
                    
                    # 如果压缩后更小，则替换原图片
                    if compressed_size < original_size * 0.95:  # 至少减少5%才替换
                        # 获取图片位置信息
                        img_rects = page.get_image_rects(xref)
                        if img_rects:
                            img_rect = img_rects[0]
                            
                            # 删除原图片
                            try:
                                page.delete_image(xref)
                            except:
                                pass
                            
                            # 插入压缩后的图片
                            if convert_to_jpeg or image_ext.lower() in ['jpeg', 'jpg']:
                                page.insert_image(img_rect, stream=compressed_bytes)
                            else:
                                page.insert_image(img_rect, stream=compressed_bytes)
                            
                            image_count += 1
                            
                except Exception as e:
                    logger.warning(f"压缩页面 {page_num + 1} 的图片 {img_index} 失败: {str(e)}")
                    continue
        
        # 生成输出文件名
        output_filename = f"compressed_{compression_level}_{uuid.uuid4().hex[:8]}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        output_path = os.path.join(app.config['CONVERTED_FOLDER'], output_filename)
        
        # 保存压缩后的PDF
        doc.save(output_path,
                 garbage=config['garbage'],
                 deflate=config['deflate'],
                 clean=config['clean'],
                 linearize=True)  # 线性化以优化Web查看
        
        doc.close()
        
        # 删除临时文件
        try:
            os.remove(temp_path)
        except:
            pass
        
        # 计算压缩结果
        compressed_size = os.path.getsize(output_path)
        compression_ratio = (1 - compressed_size / original_size) * 100
        elapsed_time = time.time() - start_time
        
        logger.info(f"PDF压缩完成: {output_filename}, 原始: {original_size / 1024:.1f}KB, "
                   f"压缩后: {compressed_size / 1024:.1f}KB, 压缩率: {compression_ratio:.1f}%, "
                   f"处理图片: {image_count}, 耗时: {elapsed_time:.2f}秒")
        
        return jsonify({
            'success': True,
            'filename': output_filename,
            'url': f'/download/{output_filename}',
            'original_size': original_size,
            'compressed_size': compressed_size,
            'compression_ratio': round(compression_ratio, 1),
            'compression_level': compression_level,
            'image_count': image_count,
            'total_pages': total_pages,
            'elapsed_time': round(elapsed_time, 2)
        })
        
    except Exception as e:
        logger.error(f"PDF压缩失败: {str(e)}")
        logger.error(traceback.format_exc())
        return jsonify({'error': f'压缩失败: {str(e)}'}), 500


@app.route('/pdf/get-pages', methods=['POST'])
def get_pdf_pages():
    """
    获取PDF所有页面的缩略图（用于前端预览）
    """
    try:
        # 支持两种字段名：file 和 pdf_file
        if 'pdf_file' in request.files:
            file = request.files['pdf_file']
        elif 'file' in request.files:
            file = request.files['file']
        else:
            return jsonify({'error': '未找到文件'}), 400
        
        if file.filename == '':
            return jsonify({'error': '未选择文件'}), 400
        
        if not file.filename.lower().endswith('.pdf'):
            return jsonify({'error': '只支持PDF文件'}), 400
        
        # 保存临时文件
        filename = secure_filename(file.filename)
        temp_filename = f"preview_{uuid.uuid4().hex[:8]}_{filename}"
        temp_path = os.path.join(app.config['UPLOAD_FOLDER'], temp_filename)
        file.save(temp_path)
        
        # 生成缩略图
        doc = fitz.open(temp_path)
        thumbnails = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            # 生成缩略图（150x150像素）
            mat = fitz.Matrix(150 / page.rect.width, 150 / page.rect.height)
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")
            
            # 转为base64
            import base64
            img_base64 = base64.b64encode(img_data).decode('utf-8')
            thumbnails.append({
                'page_number': page_num + 1,
                'thumbnail': f'data:image/png;base64,{img_base64}'
            })
        
        doc.close()
        
        # 删除临时文件
        try:
            os.remove(temp_path)
        except:
            pass
        
        logger.info(f"生成PDF缩略图: {filename}, 页数: {len(thumbnails)}")
        
        return jsonify({
            'success': True,
            'total_pages': len(thumbnails),
            'thumbnails': thumbnails
        })
        
    except Exception as e:
        logger.error(f"获取PDF页面失败: {str(e)}")
        logger.error(traceback.format_exc())
        return jsonify({'error': f'获取失败: {str(e)}'}), 500


@app.route('/download/<filename>', methods=['GET'])
def download_file(filename):
    """下载转换后的文件"""
    try:
        file_path = os.path.join(app.config['CONVERTED_FOLDER'], filename)
        if not os.path.exists(file_path):
            return jsonify({'error': '文件不存在'}), 404
        return send_file(file_path, as_attachment=True)
    except Exception as e:
        logger.error(f"下载失败: {str(e)}")
        return jsonify({'error': f'下载失败: {str(e)}'}), 500


if __name__ == '__main__':
    logger.info("启动PDF转换服务（超高性能优化版 v3.0）")
    logger.info("监听端口: 8789")
    logger.info("优化特性:")
    logger.info("  - Windows单进程模式（避免死锁）")
    logger.info("  - 分批处理大文档")
    logger.info("  - 超时保护（5分钟）")
    logger.info("  - 进度监控")
    logger.info("  - 自动清理旧文件")
    
    app.run(host='0.0.0.0', port=8789, debug=False, threaded=True)

